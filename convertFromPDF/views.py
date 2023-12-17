from django.shortcuts import render
from django.core.files.storage import FileSystemStorage
from pdf2docx import parse
from pdf2image import convert_from_path
import os
import tempfile
from pdf2image import convert_from_path
from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches
import tabula
import pandas as pd
import PyPDF2
from openpyxl import Workbook
from PIL import Image
import fitz  # PyMuPDF
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
def pdfToWord(request):
    if request.method == "POST":
        pdf_file  = request.FILES['select_file']
        fss = FileSystemStorage()
        file = fss.save(pdf_file.name, pdf_file)
        file_url = fss.url(file)
        pdf_file_name  = str(pdf_file )
        pdf_file =r'C:\Users\kushwah\OneDrive\Desktop\convert\pdfConverter\media\{}'.format(pdf_file_name)
        print(pdf_file)
        name, extension = os.path.splitext(pdf_file_name)
        print(name)
        docx_file = '{}.docx'.format(name)
        parse(pdf_file, docx_file)
        return render(request,'convertFromPdf/pdfToWord.html',{'file_url': file_url})
    return render(request,'convertFromPdf/pdfToWord.html')
    
def pdfToExcel(request):
    if request.method == "POST":
        pdf_file  = request.FILES['select_file']
        fss = FileSystemStorage()
        file = fss.save(pdf_file.name, pdf_file)
        file_url = fss.url(file)
        pdf_file_name  = str(pdf_file )
        name, extension = os.path.splitext(pdf_file_name)
        pdf_file =r'C:\Users\kushwah\OneDrive\Desktop\convert\pdfConverter\media\{}'.format(pdf_file_name)
        excel_file = r'C:\Users\kushwah\OneDrive\Desktop\convert\pdfConverter\{}.xlsx'.format(name)
        print("PdfTOExcel")
        try:
            pdf_text = ""
            wb = Workbook()
            ws = wb.active
            pdf_reader = PyPDF2.PdfReader(pdf_file)  # Use PdfReader instead
            for page in pdf_reader.pages:
                pdf_text += page.extract_text()
            pdf_lines = pdf_text.split('\n')
            for row, line in enumerate(pdf_lines):
                columns = line.split('\t')
                for col, value in enumerate(columns):
                    ws.cell(row=row+1, column=col+1, value=value)
            wb.save(excel_file)
        except Exception as e:
            print(f"Error: {str(e)}")
    return render(request,'convertFromPdf/pdfToExcel.html')


def pdfToJpg(request):
    if request.method == "POST" and 'select_file' in request.FILES:
        pdf_file  = request.FILES['select_file']
        fss = FileSystemStorage()
        file = fss.save(pdf_file.name, pdf_file)
        file_url = fss.url(file)
        pdf_file_name  = str(pdf_file )
        pdf_file =r'C:\Users\kushwah\OneDrive\Desktop\convert\pdfConverter\media\{}'.format(pdf_file_name)
        image_folder = r'C:\Users\kushwah\OneDrive\Desktop\convert\pdfConverter' 
        print("PDF TO JPG")
        try:
            pdf_document = fitz.open(pdf_file)
            for page_number in range(pdf_document.page_count):
                page = pdf_document.load_page(page_number)
        # Convert the page to an image (PNG)
                image = page.get_pixmap()
        # Create a PIL Image from the raw image data
                pil_image = Image.frombytes("RGB", [image.width, image.height], image.samples)
        # Save the PIL Image as a JPEG file
                jpeg_filename = f"{image_folder}/page_{page_number + 1}.jpg"
                pil_image.save(jpeg_filename, "JPEG")
            pdf_document.close()
        except Exception as e:
            print(f"Error: {str(e)}")
        return render(request, 'convertFromPdf/pdfToJpg.html')
    return render(request,'convertFromPdf/pdfToJpg.html')

def pdfToPowerpoint(request):
    if request.method == "POST":
        pdf_file  = request.FILES['select_file']
        fss = FileSystemStorage()
        file = fss.save(pdf_file.name, pdf_file)
        file_url = fss.url(file)
        pdf_file_name  = str(pdf_file )
        pdf_file =r'C:\Users\kushwah\OneDrive\Desktop\convert\pdfConverter\media\{}'.format(pdf_file_name)
        print(pdf_file)
        name, extension = os.path.splitext(pdf_file_name)
        pptx_file = r'C:\Users\kushwah\OneDrive\Desktop\convert\pdfConverter\output.pptx'
        try:
            prs = Presentation()
            with open(pdf_file, "rb") as pdf_file:
                pdf_reader = PdfReader(pdf_file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
            # Create a new slide
                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank layout
            # Add text to the slide
                    left = top = width = height = Inches(1)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    text_frame.text = text
            # Save the PowerPoint presentation
            prs.save(pptx_file)
        except Exception as e:
            print(f"Error: {str(e)}")    
        return render(request,'convertFromPdf/pdfToPowerponit.html')
    return render(request,'convertFromPdf/pdfToPowerponit.html')